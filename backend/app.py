# app.py
import os
import json
import time
import tempfile
import base64
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import PyPDF2
from docx import Document
from pptx import Presentation
import openai
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.units import inch
from gtts import gTTS
import uuid
import threading
import shutil
from langchain.text_splitter import RecursiveCharacterTextSplitter
from openai import OpenAI

# Initialize Flask application
app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# Configure OpenAI API key - replace with your key
client = OpenAI(api_key="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")  # Replace this with your actual OpenAI API key

# Create upload directory if it doesn't exist
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'output'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Define allowed file extensions
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'ppt', 'pptx'}

# Helper function to check allowed file extensions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

# Function to extract text from PDF
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = PyPDF2.PdfReader(file)
        num_pages = len(pdf_reader.pages)
        for page_num in range(num_pages):
            page = pdf_reader.pages[page_num]
            text += page.extract_text() + "\n"
    return text

# Function to extract text from DOCX
def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# Function to extract text from PPT/PPTX
def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to process document and extract text
def process_document(file_path):
    file_extension = file_path.split('.')[-1].lower()
    if file_extension == 'pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == 'docx':
        return extract_text_from_docx(file_path)
    elif file_extension in ['ppt', 'pptx']:
        return extract_text_from_ppt(file_path)
    else:
        return None

# Split large text into chunks for processing
def split_text(text, chunk_size=15000, chunk_overlap=500):
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

# Create PDF with text and audio
def create_pdf_with_audio(content, title, session_id, output_path):
    # Create a unique filename
    filename = f"{title.replace(' ', '_')}_{session_id}.pdf"
    filepath = os.path.join(output_path, filename)
    
    # Create PDF
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    # Add title
    title_style = ParagraphStyle(
        'Title',
        parent=styles['Heading1'],
        fontSize=16,
        spaceAfter=12
    )
    story.append(Paragraph(title, title_style))
    story.append(Spacer(1, 0.2*inch))
    
    # Add content
    content_style = ParagraphStyle(
        'Content',
        parent=styles['Normal'],
        fontSize=12,
        spaceAfter=10
    )
    
    paragraphs = content.split('\n')
    for paragraph in paragraphs:
        if paragraph.strip():
            story.append(Paragraph(paragraph, content_style))
            story.append(Spacer(1, 0.1*inch))
    
    # Build PDF
    doc.build(story)
    
    # Create audio file
    audio_filename = f"{title.replace(' ', '_')}_{session_id}.mp3"
    audio_filepath = os.path.join(output_path, audio_filename)
    
    # Break content into smaller chunks for TTS (to avoid gTTS limitations)
    max_tts_length = 5000
    content_chunks = [content[i:i+max_tts_length] for i in range(0, len(content), max_tts_length)]
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_audio_files = []
        
        # Create audio chunks
        for i, chunk in enumerate(content_chunks):
            if not chunk.strip():
                continue
                
            temp_file = os.path.join(temp_dir, f"chunk_{i}.mp3")
            tts = gTTS(text=chunk, lang='en', slow=False)
            tts.save(temp_file)
            temp_audio_files.append(temp_file)
        
        # Combine audio chunks if there are multiple
        if len(temp_audio_files) > 1:
            # Here you would need to use a library like pydub to concatenate the audio files
            # For simplicity, we'll just use the first chunk
            shutil.copyfile(temp_audio_files[0], audio_filepath)
        elif len(temp_audio_files) == 1:
            shutil.copyfile(temp_audio_files[0], audio_filepath)
    
    return filename, audio_filename

# Generate study materials using OpenAI
def generate_study_materials(document_text, num_pages, session_id):
    output_path = os.path.join(OUTPUT_FOLDER, session_id)
    os.makedirs(output_path, exist_ok=True)
    
    # Calculate appropriate sizes for materials based on document length
    summary_pages = max(min(num_pages // 10, 50), 5)  # Between 5 and 50 pages
    qa_count = max(min(num_pages * 2, 100), 20)  # Between 20 and 100 questions
    quiz_count = max(min(num_pages, 50), 10)  # Between 10 and 50 quiz questions
    
    # Break document into manageable chunks
    chunks = split_text(document_text)
    
    # Dictionary to store the results
    results = {
        "summary": "",
        "qa": "",
        "quiz": "",
        "important_elements": ""
    }
    
    # Process each chunk and accumulate results
    for i, chunk in enumerate(chunks):
        print(f"Processing chunk {i+1} of {len(chunks)}")
        
        # Generate summary
        summary_prompt = f"""
        Create a comprehensive exam-focused study summary based on the following document excerpt:
        
        Text: {chunk}
        
        Guidelines:
        - Focus exclusively on the technical content, concepts, definitions, and examples from the document
        - Omit any references to authors, universities, or other publication details 
        - Organize information in a structured, easily digestible format for exam preparation
        - Include only information that would be relevant for assessment purposes
        - Highlight key terms, definitions, relationships between concepts, and important distinctions
        - The full document is approximately {num_pages} pages, so ensure your summary is proportionally detailed
        - Use bullet points, numbered lists, or other formatting to enhance readability when appropriate
        
        Your summary should serve as a comprehensive study aid that captures all exam-relevant information from this portion of the document.
        """
        summary_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": summary_prompt}],
            max_tokens=2000,
            temperature=0.7
        )
        results["summary"] += summary_response.choices[0].message.content
        
        # Generate Q&A
        qa_prompt = f"""
        Create comprehensive questions and answers based on the following text:
        
        Text: {chunk}
        
        Generate two distinct question sets - one with short-answer questions and one with long-answer questions. Include a variety of questions that test factual recall, conceptual understanding, and application of knowledge.
        
        Format your response with all short-answer questions grouped together first, followed by all long-answer questions:
        
        SHORT ANSWER QUESTIONS:
        Q1: [Question]
        A1: [Concise answer - typically 1-3 sentences]
        
        Q2: [Question]
        A2: [Concise answer - typically 1-3 sentences]
        
        [Continue with all short-answer questions]
        
        LONG ANSWER QUESTIONS:
        Q1: [Question]
        A1: [Comprehensive answer spanning multiple paragraphs that thoroughly explains the concept]
        
        Q2: [Question]
        A2: [Comprehensive answer spanning multiple paragraphs that thoroughly explains the concept]
        
        [Continue with all long-answer questions]
        
        Note: Ensure all short-answer questions are grouped together in the first section, and all long-answer questions are grouped together in the second section. Do not alternate between short and long answers.
        """
        qa_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": qa_prompt}],
            max_tokens=2000,
            temperature=0.7
        )
        results["qa"] += qa_response.choices[0].message.content + "\n\n"
        
        # Generate quiz
        quiz_prompt = f"""
        Create multiple-choice quiz questions with answers based on the following text:
        
        Text: {chunk}
        
        Format as:
        Q1: [Question]
        A) [Option A]
        B) [Option B]
        C) [Option C]
        D) [Option D]
        Answer: [Correct option letter]
        Explanation: [Why this is the correct answer]
        """
        quiz_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": quiz_prompt}],
            max_tokens=1500,
            temperature=0.7
        )
        results["quiz"] += quiz_response.choices[0].message.content + "\n\n"
        
        # Generate important elements
        elements_prompt = f"""
        Extract and organize the most important elements from the following text:
        
        Text: {chunk}
        
        Format your response into three consolidated sections:
        
        EQUATIONS & FORMULAS:
        - [List all equations and formulas with brief explanations]
        - [Include mathematical notations if present]
        
        KEY CONCEPTS:
        - [List all important concepts with concise explanations]
        - [Group related concepts together]
        
        DIAGRAMS & FLOWCHARTS:
        - [Detailed descriptions of any diagrams or flowcharts]
        - [Explain what each visual represents, its components, and significance]
        - [If a diagram is mentioned but not fully described, note what information would be needed]
        
        Note: For each section, combine all related items rather than repeating section headers. Present information in a logical, organized manner.
        """
        elements_response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": elements_prompt}],
            max_tokens=1500,
            temperature=0.7
        )
        results["important_elements"] += elements_response.choices[0].message.content + "\n\n"
    
    # Create PDFs with audio for each type of content
    output_files = {}
    
    summary_pdf, summary_audio = create_pdf_with_audio(
        results["summary"], 
        "Comprehensive Summary", 
        session_id,
        output_path
    )
    output_files["summary"] = {"pdf": summary_pdf, "audio": summary_audio}
    
    qa_pdf, qa_audio = create_pdf_with_audio(
        results["qa"], 
        "Questions and Answers", 
        session_id,
        output_path
    )
    output_files["qa"] = {"pdf": qa_pdf, "audio": qa_audio}
    
    quiz_pdf, quiz_audio = create_pdf_with_audio(
        results["quiz"], 
        "Practice Quiz", 
        session_id,
        output_path
    )
    output_files["quiz"] = {"pdf": quiz_pdf, "audio": quiz_audio}
    
    elements_pdf, elements_audio = create_pdf_with_audio(
        results["important_elements"], 
        "Important Elements", 
        session_id,
        output_path
    )
    output_files["important_elements"] = {"pdf": elements_pdf, "audio": elements_audio}
    
    return output_files

# Global dictionary to store processing status
processing_status = {}

# Process the document in a separate thread
def process_document_async(file_path, original_filename, session_id):
    try:
        # Update status to "processing"
        processing_status[session_id] = {
            "status": "processing",
            "message": "Extracting text from document..."
        }
        
        # Extract text from the document
        document_text = process_document(file_path)
        if not document_text:
            processing_status[session_id] = {
                "status": "error",
                "message": "Failed to extract text from the document."
            }
            return
        
        # Estimate number of pages based on characters
        # This is a rough estimate; actual pages depend on formatting
        chars_per_page = 2500  # Average characters per page
        estimated_pages = max(1, len(document_text) // chars_per_page)
        
        processing_status[session_id] = {
            "status": "processing",
            "message": f"Text extracted. Generating study materials for approximately {estimated_pages} pages..."
        }
        
        # Generate study materials
        output_files = generate_study_materials(document_text, estimated_pages, session_id)
        
        # Update status to "completed"
        processing_status[session_id] = {
            "status": "completed",
            "message": "Processing completed successfully.",
            "files": output_files
        }
    except Exception as e:
        processing_status[session_id] = {
            "status": "error",
            "message": f"An error occurred: {str(e)}"
        }

# Route for uploading documents
@app.route('/upload', methods=['POST'])
def upload_file():
    # Check if the post request has the file part
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    
    # If user submits empty form
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if file and allowed_file(file.filename):
        # Generate a unique session ID
        session_id = str(uuid.uuid4())
        
        # Create session directory
        session_dir = os.path.join(UPLOAD_FOLDER, session_id)
        os.makedirs(session_dir, exist_ok=True)
        
        # Save file
        filename = secure_filename(file.filename)
        file_path = os.path.join(session_dir, filename)
        file.save(file_path)
        
        # Start processing in a separate thread
        processing_thread = threading.Thread(
            target=process_document_async,
            args=(file_path, filename, session_id)
        )
        processing_thread.start()
        
        return jsonify({
            'message': 'File uploaded successfully',
            'session_id': session_id
        }), 200
    
    return jsonify({'error': 'File type not allowed'}), 400

# Route for checking processing status
@app.route('/status/<session_id>', methods=['GET'])
def check_status(session_id):
    if session_id in processing_status:
        return jsonify(processing_status[session_id]), 200
    else:
        return jsonify({'status': 'not_found', 'message': 'Session ID not found'}), 404

# Route for downloading generated files
@app.route('/download/<session_id>/<file_type>/<file_name>', methods=['GET'])
def download_file(session_id, file_type, file_name):
    file_path = os.path.join(OUTPUT_FOLDER, session_id, file_name)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    else:
        return jsonify({'error': 'File not found'}), 404

# Main entry point
if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)